#!/usr/bin/env python3
"""从 PPTX 中提取原始 shape 几何信息、文本和字体元数据。

输出单位：
- 坐标与尺寸使用 inch
- 字号使用 pt
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any, Dict, Iterable, List, Tuple


def localize_argparse() -> None:
    """将 argparse 自动生成的帮助/错误文案改为中文。"""
    translations = {
        "usage: ": "用法：",
        "positional arguments": "位置参数",
        "options": "选项",
        "optional arguments": "可选参数",
        "show this help message and exit": "显示此帮助信息并退出",
        "the following arguments are required: %s": "缺少必填参数：%s",
        "unrecognized arguments: %s": "无法识别的参数：%s",
        "invalid choice: %(value)r (choose from %(choices)s)": "无效选择：%(value)r（可选：%(choices)s）",
        "argument %(argument_name)s: %(message)s": "参数 %(argument_name)s：%(message)s",
        "%(prog)s: error: %(message)s\n": "%(prog)s：错误：%(message)s\n",
        "expected one argument": "需要一个参数",
        "not allowed with argument %s": "不能与参数 %s 同时使用",
        "ignored explicit argument %r": "忽略显式参数 %r",
    }
    argparse._ = lambda text: translations.get(text, text)  # type: ignore[attr-defined]



EMU_PER_INCH = 914_400


def length_to_inches(value: Any) -> float | None:
    if value is None:
        return None
    try:
        return round(float(value) / EMU_PER_INCH, 4)
    except Exception:
        return None


def font_size_to_points(value: Any) -> float | None:
    if value is None:
        return None
    try:
        # python-pptx 的 Length 对象通常提供 .pt 属性。
        return round(float(value.pt), 3)
    except Exception:
        try:
            return round(float(value) / 12_700, 3)
        except Exception:
            return None


def safe_enum(value: Any) -> str | None:
    if value is None:
        return None
    try:
        return str(value)
    except Exception:
        return None


def safe_int(value: Any) -> int | None:
    try:
        return int(value)
    except Exception:
        return None


def color_to_value(font: Any) -> str | None:
    try:
        color = font.color
        if color is None:
            return None
        # 适用于显式 RGB 颜色；主题色或继承色可能无法直接读取。
        rgb = color.rgb
        return str(rgb) if rgb is not None else None
    except Exception:
        return None


def font_to_dict(font: Any) -> Dict[str, Any]:
    if font is None:
        return {}
    return {
        "name": getattr(font, "name", None),
        "size_pt": font_size_to_points(getattr(font, "size", None)),
        "bold": getattr(font, "bold", None),
        "italic": getattr(font, "italic", None),
        "underline": getattr(font, "underline", None),
        "color_rgb": color_to_value(font),
    }


def text_frame_to_dict(shape: Any) -> Dict[str, Any] | None:
    if not getattr(shape, "has_text_frame", False):
        return None

    tf = shape.text_frame
    paragraphs: List[Dict[str, Any]] = []
    has_runs = False

    for pi, paragraph in enumerate(tf.paragraphs):
        runs: List[Dict[str, Any]] = []
        for ri, run in enumerate(paragraph.runs):
            has_runs = True
            runs.append(
                {
                    "run_index": ri,
                    "text": run.text,
                    "font": font_to_dict(run.font),
                }
            )

        paragraphs.append(
            {
                "paragraph_index": pi,
                "text": paragraph.text,
                "level": getattr(paragraph, "level", None),
                "alignment": safe_enum(getattr(paragraph, "alignment", None)),
                "font": font_to_dict(getattr(paragraph, "font", None)),
                "runs": runs,
            }
        )

    text = getattr(shape, "text", "")
    return {
        "text": text,
        "has_runs": has_runs,
        "paragraphs": paragraphs,
        "margin_left_in": length_to_inches(getattr(tf, "margin_left", None)),
        "margin_right_in": length_to_inches(getattr(tf, "margin_right", None)),
        "margin_top_in": length_to_inches(getattr(tf, "margin_top", None)),
        "margin_bottom_in": length_to_inches(getattr(tf, "margin_bottom", None)),
        "word_wrap": getattr(tf, "word_wrap", None),
        "auto_size": safe_enum(getattr(tf, "auto_size", None)),
    }


def iter_shapes(shapes: Iterable[Any], group_path: Tuple[int, ...] = ()) -> Iterable[Tuple[Any, Tuple[int, ...]]]:
    for idx, shape in enumerate(shapes):
        path = group_path + (idx,)
        # 组合形状有 .shapes 集合；同时产出组合本身和内部子元素。
        yield shape, path
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            yield from iter_shapes(nested, path)


def shape_to_dict(shape: Any, slide_index: int, path: Tuple[int, ...]) -> Dict[str, Any]:
    left = length_to_inches(getattr(shape, "left", None))
    top = length_to_inches(getattr(shape, "top", None))
    width = length_to_inches(getattr(shape, "width", None))
    height = length_to_inches(getattr(shape, "height", None))
    bbox = [left, top, width, height]

    center = None
    if None not in bbox:
        center = [round(left + width / 2, 4), round(top + height / 2, 4)]  # type: ignore[operator]

    xml = ""
    try:
        xml = shape.element.xml
    except Exception:
        pass

    text_frame = text_frame_to_dict(shape)

    return {
        "slide_index": slide_index,
        "shape_id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", None),
        "path": list(path),
        "shape_type_id": safe_int(getattr(shape, "shape_type", None)),
        "shape_type": safe_enum(getattr(shape, "shape_type", None)),
        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
        "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
        "has_table": bool(getattr(shape, "has_table", False)),
        "has_chart": bool(getattr(shape, "has_chart", False)),
        "has_picture": safe_int(getattr(shape, "shape_type", None)) == 13,
        "bbox": bbox,
        "center": center,
        "text_frame": text_frame,
        "contains_omml_xml": ("<m:oMath" in xml) or ("<m:oMathPara" in xml),
    }


def extract(pptx_path: Path, include_images: bool = False) -> Dict[str, Any]:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - dependency error path
        raise SystemExit(
            "缺少依赖 python-pptx，请先安装：python3 -m pip install python-pptx"
        ) from exc

    prs = Presentation(str(pptx_path))
    slides: List[Dict[str, Any]] = []

    for slide_index, slide in enumerate(prs.slides):
        slide_shapes: List[Dict[str, Any]] = []
        for shape, path in iter_shapes(slide.shapes):
            shape_type_id = safe_int(getattr(shape, "shape_type", None))
            if shape_type_id == 13 and not include_images:
                continue
            slide_shapes.append(shape_to_dict(shape, slide_index, path))

        slides.append(
            {
                "slide_index": slide_index,
                "shape_count": len(slide_shapes),
                "shapes": slide_shapes,
            }
        )

    return {
        "source_pptx": str(pptx_path),
        "slide_width_in": length_to_inches(prs.slide_width),
        "slide_height_in": length_to_inches(prs.slide_height),
        "slide_count": len(slides),
        "slides": slides,
    }


def main(argv: List[str] | None = None) -> int:
    localize_argparse()
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path, help="编辑后的 PPTX 或新鲜渲染的 PPTX 路径")
    parser.add_argument("--output", "-o", type=Path, default=Path("pptx_raw_shapes.json"), help="输出原始 shape JSON 的路径")
    parser.add_argument("--include-images", action="store_true", help="包含图片/背景图 shape")
    args = parser.parse_args(argv)

    if not args.pptx.exists():
        parser.error(f"找不到 PPTX 文件：{args.pptx}")

    result = extract(args.pptx, include_images=args.include_images)
    args.output.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")

    text_shapes = sum(
        1
        for slide in result["slides"]
        for shape in slide["shapes"]
        if shape.get("has_text_frame")
    )
    print(
        f"已写入 {args.output} —— 共 {result['slide_count']} 页，"
        f"提取 {text_shapes} 个文本 shape。",
        file=sys.stderr,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
