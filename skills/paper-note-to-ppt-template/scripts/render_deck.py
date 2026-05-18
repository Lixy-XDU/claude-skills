#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将 layout_plan.json 渲染成 PPTX。

设计原则：
- PNG 模板作为整页背景。
- renderer 只在预设 slot 内叠加透明文本 / 图片 / 公式。
- 不额外绘制卡片背景，避免和模板卡片冲突。
"""

from __future__ import annotations

import json
import sys
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from PIL import Image
except Exception:  # pragma: no cover
    Image = None

SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPT_DIR))
try:
    from render_formula import render_formula
except Exception:
    render_formula = None

# 累加每页公式渲染状态，build() 末尾写入 formula_report.json
_formula_results: list[dict] = []


def inch_box(bbox: list[float]):
    return [Inches(float(v)) for v in bbox]


def hex_to_rgb(value: str) -> RGBColor:
    value = value.strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def get_style(theme: dict, style_name: str) -> dict:
    return theme.get("styles", {}).get(style_name, theme.get("styles", {}).get("body", {}))


def resolve_color(theme: dict, color_name: str | None) -> RGBColor:
    colors = theme.get("colors", {})
    value = colors.get(color_name or "text_primary", color_name or "#082B5F")
    return hex_to_rgb(value)


def set_run_style(run, theme: dict, style: dict):
    font = run.font
    font.size = Pt(float(style.get("font_size", 14)))
    font.bold = bool(style.get("bold", False))
    font.name = style.get("font_face", theme.get("fonts", {}).get("zh", "Microsoft YaHei"))
    font.color.rgb = resolve_color(theme, style.get("color", "text_primary"))


def add_textbox(slide, text: str, bbox: list[float], theme: dict, style_name: str, align: str = "left"):
    if text is None:
        return None
    text = str(text).strip()
    if not text:
        return None
    x, y, w, h = inch_box(bbox)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = ""
    p.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    set_run_style(run, theme, get_style(theme, style_name))
    return box


def add_multiline_text(slide, lines: list[str], bbox: list[float], theme: dict, style_name: str, bullet: bool = False):
    x, y, w, h = inch_box(bbox)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    style = get_style(theme, style_name)
    for i, line in enumerate([str(v).strip() for v in lines if str(v).strip()]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = ("• " if bullet else "") + line
        set_run_style(run, theme, style)
    return box


def fit_image_size(image_path: Path, bbox: list[float]) -> tuple[int, int, int, int]:
    x, y, w, h = inch_box(bbox)
    if Image is None:
        return x, y, w, h
    try:
        with Image.open(image_path) as im:
            iw, ih = im.size
        box_ratio = float(w) / float(h)
        img_ratio = iw / ih
        if img_ratio > box_ratio:
            new_w = w
            new_h = int(float(w) / img_ratio)
        else:
            new_h = h
            new_w = int(float(h) * img_ratio)
        new_x = x + int((float(w) - float(new_w)) / 2)
        new_y = y + int((float(h) - float(new_h)) / 2)
        return new_x, new_y, new_w, new_h
    except Exception:
        return x, y, w, h


def add_picture_contain(slide, image_path: Path, bbox: list[float]):
    if not image_path.exists():
        # 画一个极轻提示，便于发现问题
        add_textbox(slide, f"[缺少图片: {image_path.name}]", bbox, {"styles": {"body": {"font_size": 12, "color": "#64748B"}}, "colors": {}}, "body")
        return None
    x, y, w, h = fit_image_size(image_path, bbox)
    return slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)


def _set_pic_opacity(pic, opacity_pct: int):
    """设置图片透明度。opacity_pct: 0-100, 100=完全不透明, 0=完全透明。"""
    try:
        from lxml import etree as lx
    except ImportError:
        return
    if not isinstance(pic, tuple):
        pic = (pic,)
    for p in pic:
        blip = p._element.find('.//' + _qn('a:blip'))
        if blip is None:
            continue
        # 移除已有的 alphaModFix
        for old in blip.findall(_qn('a:alphaModFix')):
            blip.remove(old)
        alpha = lx.SubElement(blip, _qn('a:alphaModFix'))
        alpha.set('amt', str(opacity_pct * 1000))  # 80% → 80000


def _qn(tag: str) -> str:
    """将 a:xxx 或 p:xxx 前缀标签转为命名空间限定的 XPath 查询字符串。"""
    from pptx.oxml.ns import qn as _pptx_qn
    return _pptx_qn(tag)


def find_template(template_dir: Path, aliases: list[str]) -> Path | None:
    for name in aliases:
        p = template_dir / name
        if p.exists():
            return p
    stems = {Path(a).stem.lower() for a in aliases}
    for p in template_dir.glob("*"):
        if p.is_file() and p.suffix.lower() in {".png", ".jpg", ".jpeg"}:
            if p.stem.lower() in stems:
                return p
    return None


def resolve_path(value: str, plan_dir: Path, output_dir: Path) -> Path:
    p = Path(value).expanduser()
    if p.is_absolute():
        return p
    candidates = [plan_dir / p, output_dir / p, Path.cwd() / p]
    for c in candidates:
        if c.exists():
            return c.resolve()
    return (plan_dir / p).resolve()


def plain_text_from_mixed(value: Any) -> tuple[str, str, list[str]]:
    if isinstance(value, str):
        return "", value, []
    if isinstance(value, dict):
        title = str(value.get("title", "")).strip()
        body = str(value.get("body", "")).strip()
        bullets = value.get("bullets", value.get("items", [])) or []
        if isinstance(bullets, str):
            bullets = [bullets]
        return title, body, [str(x).strip() for x in bullets if str(x).strip()]
    return "", str(value), []


def render_mixed_box(slide, value: Any, bbox: list[float], theme: dict, plan_dir: Path, output_dir: Path, slide_no: int, slot_name: str):
    # 支持 image_path / latex / title+body+bullets
    if isinstance(value, dict) and value.get("image_path"):
        img = resolve_path(str(value["image_path"]), plan_dir, output_dir)
        cap = str(value.get("caption", "")).strip()
        image_box = [bbox[0] + 0.12, bbox[1] + 0.15, bbox[2] - 0.24, bbox[3] - (0.55 if cap else 0.3)]
        add_picture_contain(slide, img, image_box)
        if cap:
            add_textbox(slide, cap, [bbox[0] + 0.18, bbox[1] + bbox[3] - 0.35, bbox[2] - 0.36, 0.25], theme, "caption", align="center")
        return

    if isinstance(value, dict) and value.get("latex"):
        # 公式不允许放在 mixed_box 中，记录警告并跳过
        _formula_results.append({
            "slide_no": slide_no,
            "slot": slot_name,
            "latex": str(value.get("latex", ""))[:120],
            "status": "skipped_in_mixed_box",
            "message": "LaTeX found in mixed_box slot; formulas must use formula slots only. Skipped.",
        })
        return

    title, body, bullets = plain_text_from_mixed(value)
    cur_y = bbox[1] + 0.0
    if title:
        add_textbox(slide, title, [bbox[0] + 0.25, cur_y, bbox[2] - 0.5, 0.44], theme, "box_title")
        cur_y += 0.70
    if body:
        body_h = 1.15 if bullets else max(0.6, bbox[3] - (cur_y - bbox[1]) - 0.3)
        add_textbox(slide, body, [bbox[0] + 0.25, cur_y, bbox[2] - 0.5, body_h], theme, "box_body")
        cur_y += min(body_h, 1.2) + 0.12
    if bullets:
        add_multiline_text(slide, bullets[:5], [bbox[0] + 0.25, cur_y, bbox[2] - 0.5, bbox[1] + bbox[3] - cur_y - 0.15], theme, "bullet", bullet=True)


def render_agenda(slide, items: list[str], bboxes: list[list[float]], theme: dict):
    for item, bbox in zip(items[:len(bboxes)], bboxes):
        add_textbox(slide, str(item), bbox, theme, "agenda")


def render_inline_meta(slide, items: list[str], bbox: list[float], theme: dict):
    if not items:
        return
    text = "    |    ".join(str(x).strip() for x in items if str(x).strip())
    add_textbox(slide, text, bbox, theme, "meta")


def render_slide(prs: Presentation, slide_plan: dict, layout_spec: dict, template_path: Path, theme: dict, plan_dir: Path, output_dir: Path):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # 背景铺满，80% 透明度，图片以嵌入方式存储
    bg_pic = slide.shapes.add_picture(str(template_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    _set_pic_opacity(bg_pic, 80)

    slide_no = int(slide_plan.get("slide_no", len(prs.slides)))
    slots = dict(slide_plan.get("slots", {}) or {})
    # 让 slide_plan.title 可自动填入 title slot
    if "title" in layout_spec.get("slots", {}) and "title" not in slots and slide_plan.get("title"):
        slots["title"] = slide_plan.get("title")

    for slot_name, spec in layout_spec.get("slots", {}).items():
        value = slots.get(slot_name)
        if value is None or value == "":
            continue
        kind = spec.get("kind")
        bbox = spec.get("bbox")
        style = spec.get("style", "body")
        if kind == "text":
            align = "center" if slot_name in {"title", "subtitle"} else "left"
            add_textbox(slide, str(value), bbox, theme, style, align=align)
        elif kind == "inline_meta":
            if isinstance(value, str):
                value = [value]
            render_inline_meta(slide, value, bbox, theme)
        elif kind == "agenda_items":
            if isinstance(value, dict):
                value = value.get("items", [])
            render_agenda(slide, [str(x) for x in value], spec.get("bboxes", []), theme)
        elif kind == "bullets":
            add_multiline_text(slide, [str(x) for x in value], bbox, theme, style, bullet=True)
        elif kind == "image":
            if isinstance(value, dict):
                image_path = value.get("image_path") or value.get("path")
                cap = str(value.get("caption", "")).strip()
            else:
                image_path = str(value)
                cap = ""
            if cap:
                cap_h = 0.35
                image_bbox = [bbox[0], bbox[1], bbox[2], bbox[3] - cap_h]
                if image_path:
                    add_picture_contain(slide, resolve_path(str(image_path), plan_dir, output_dir), image_bbox)
                add_textbox(slide, cap, [bbox[0], bbox[1] + bbox[3] - cap_h, bbox[2], cap_h], theme, "caption", align="center")
            else:
                if image_path:
                    add_picture_contain(slide, resolve_path(str(image_path), plan_dir, output_dir), bbox)
        elif kind == "formula":
            if isinstance(value, dict) and value.get("latex"):
                latex = str(value.get("latex", "")).strip()
                formula_dir = output_dir / "assets" / "formulas"
                formula_dir.mkdir(parents=True, exist_ok=True)
                out = formula_dir / f"slide_{slide_no:02d}_{slot_name}.png"
                if render_formula:
                    render_formula(latex, out)
                add_picture_contain(slide, out, bbox)
                _formula_results.append({
                    "slide_no": slide_no,
                    "slot": slot_name,
                    "latex": latex[:120],
                    "status": "png_rendered",
                    "message": "Formula rendered as PNG image.",
                })
            elif isinstance(value, dict) and value.get("image_path"):
                add_picture_contain(slide, resolve_path(str(value["image_path"]), plan_dir, output_dir), bbox)
                _formula_results.append({
                    "slide_no": slide_no,
                    "slot": slot_name,
                    "latex": "",
                    "status": "image_path",
                    "message": "Pre-rendered formula image used.",
                })
        elif kind == "mixed_box":
            render_mixed_box(slide, value, bbox, theme, plan_dir, output_dir, slide_no, slot_name)
        elif kind == "caption":
            add_textbox(slide, str(value), bbox, theme, style, align="center")
        else:
            # fallback text
            if bbox:
                add_textbox(slide, str(value), bbox, theme, style)

    notes = str(slide_plan.get("speaker_notes", "")).strip()
    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass


def build(layout_plan_path: Path, layouts_path: Path, theme_path: Path, template_dir: Path, output_path: Path) -> None:
    global _formula_results
    _formula_results.clear()

    layout_plan = json.loads(layout_plan_path.read_text(encoding="utf-8"))
    layouts_data = json.loads(layouts_path.read_text(encoding="utf-8"))
    theme = json.loads(theme_path.read_text(encoding="utf-8"))
    template_aliases = layouts_data.get("template_aliases", {})
    layouts = layouts_data.get("layouts", {})

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(float(layouts_data.get("canvas", {}).get("width", 13.333)))
    prs.slide_height = Inches(float(layouts_data.get("canvas", {}).get("height", 7.5)))

    plan_dir = layout_plan_path.parent.resolve()
    output_dir = output_path.parent.resolve()

    missing = []
    for slide in layout_plan.get("slides", []):
        layout_name = slide.get("layout")
        if layout_name not in layouts:
            raise ValueError(f"Unknown layout: {layout_name}")
        layout_spec = layouts[layout_name]
        template_key = layout_spec.get("template_key", layout_name)
        template_path = find_template(template_dir, template_aliases.get(template_key, []))
        if not template_path:
            missing.append(template_key)
            continue
        render_slide(prs, slide, layout_spec, template_path, theme, plan_dir, output_dir)

    if missing:
        raise FileNotFoundError("Missing template images: " + ", ".join(sorted(set(missing))))

    prs.save(output_path)

    # 写公式渲染报告
    formula_report_path = output_path.parent / "formula_report.json"
    formula_report_path.write_text(
        json.dumps({"formula_results": _formula_results}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def main() -> int:
    if len(sys.argv) != 6:
        print("Usage: render_deck.py <layout_plan.json> <template_layouts.json> <theme_tokens.json> <template_dir> <output.pptx>")
        return 2
    build(Path(sys.argv[1]), Path(sys.argv[2]), Path(sys.argv[3]), Path(sys.argv[4]), Path(sys.argv[5]))
    print(f"Wrote {sys.argv[5]}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
