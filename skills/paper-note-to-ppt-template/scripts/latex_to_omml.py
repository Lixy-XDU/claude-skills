#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
LaTeX → MathML → OMML → PowerPoint 原生公式管道。

使用 latex2mathml + Microsoft MML2OMML.xsl (XSLT) 转换。

函数：
  pre_flight_check()        → 检查 latex2mathml / lxml / MML2OMML.xsl 是否可用
  latex_to_omml(latex)      → LaTeX → OMML lxml Element
  insert_native_formula_omml(slide, latex, bbox, ...) → 完整的插入流程
"""

from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Any

SCRIPT_DIR = Path(__file__).resolve().parent
MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"

# ── XSLT 缓存 ───────────────────────────────────────────────────────

_xslt_transform = None


def _get_xslt():
    """加载 MML2OMML.xsl XSLT 转换器（缓存）。"""
    global _xslt_transform
    if _xslt_transform is not None:
        return _xslt_transform

    import lxml.etree as LX

    # 在常见 Office 安装位置寻找 MML2OMML.xsl
    candidates = [
        Path("C:/Program Files/Microsoft Office/root/Office16/MML2OMML.XSL"),
        Path("C:/Program Files (x86)/Microsoft Office/root/Office16/MML2OMML.XSL"),
        SCRIPT_DIR / "MML2OMML.XSL",
    ]
    xsl_path = None
    for c in candidates:
        if c.exists():
            xsl_path = c
            break
    if xsl_path is None:
        raise FileNotFoundError(
            "MML2OMML.xsl not found. Place it in the scripts/ directory, "
            "or install Microsoft Office."
        )

    xslt_doc = LX.parse(str(xsl_path))
    _xslt_transform = LX.XSLT(xslt_doc)
    return _xslt_transform


# ── public API ───────────────────────────────────────────────────────


def pre_flight_check() -> dict:
    """检查所有依赖是否可用。"""
    result: dict[str, Any] = {"available": True, "message": "ok", "checks": {}}

    # latex2mathml
    try:
        from latex2mathml.converter import convert  # noqa: F401
        result["checks"]["latex2mathml"] = True
    except ImportError:
        result["checks"]["latex2mathml"] = False
        result["available"] = False
        result["message"] = "latex2mathml not installed. pip install latex2mathml"

    # lxml
    try:
        import lxml.etree  # noqa: F401
        result["checks"]["lxml"] = True
    except ImportError:
        result["checks"]["lxml"] = False
        result["available"] = False
        result["message"] = "lxml not installed. pip install lxml"

    # MML2OMML.xsl
    try:
        _get_xslt()
        result["checks"]["MML2OMML.xsl"] = True
    except Exception as exc:
        result["checks"]["MML2OMML.xsl"] = False
        result["available"] = False
        result["message"] = str(exc)

    return result


def latex_to_omml(latex: str) -> "lxml.etree.Element":
    """Convert LaTeX → MathML → OMML via Microsoft XSLT, return lxml Element.

    Returns <m:oMathPara> containing <m:oMath> with the converted formula.
    """
    import lxml.etree as LX
    from latex2mathml.converter import convert as latex2mathml_convert

    # Step 1: LaTeX → MathML
    mathml_str = latex2mathml_convert(latex.strip())

    # Step 2: MathML → OMML via Microsoft XSLT
    xslt = _get_xslt()
    mathml_doc = LX.fromstring(mathml_str.encode("utf-8"))
    omml_result = xslt(mathml_doc)

    # XSLT 输出 <m:oMath> 元素，用 deepcopy 取出，
    # 挂到 <m:oMathPara> 下（不走字符串序列化，避免结构损坏）
    from copy import deepcopy
    oMath_elem = deepcopy(omml_result.getroot())

    oMathPara = LX.Element(f"{{{MATH_NS}}}oMathPara", nsmap={"m": MATH_NS})
    oMathPara.append(oMath_elem)

    return oMathPara


def insert_native_formula_omml(
    slide,
    latex: str,
    bbox: list[float],
    output_dir: Path,
    slide_no: int,
    slot_name: str,
    fallback: str = "fail_then_report",
) -> dict:
    """将 LaTeX 公式以原生 OMML 插入 PPT 幻灯片。

    返回 {"status": "ok"|"fallback_image"|"failed", "message": "...", "image_path": "..."}
    """
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn

    try:
        omml_root = latex_to_omml(latex)
    except Exception as exc:
        if fallback == "image":
            return _fallback_image(latex, slide, bbox, output_dir, slide_no, slot_name, str(exc))
        return {
            "status": "failed",
            "message": f"LaTeX -> OMML conversion failed: {exc}",
        }

    try:
        x, y, w, h = [Inches(float(v)) for v in bbox]
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 禁用 auto-fit，防止公式框因无可见文本被收缩到 0
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            for af in list(bodyPr):
                tag = af.tag.split("}", 1)[1] if "}" in af.tag else af.tag
                if tag in ("spAutoFit", "noAutofit", "normAutofit"):
                    bodyPr.remove(af)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        # 删除 python-pptx 自动创建的空 <a:r> 和 <a:pPr>，
        # 然后将 OMML 作为 <a:p> 的直接子节点插入（符合 OOXML 规范）。
        for r_elem in p._p.findall(qn("a:r")):
            p._p.remove(r_elem)
        for pp in p._p.findall(qn("a:pPr")):
            p._p.remove(pp)
        p._p.append(omml_root)

        return {
            "status": "ok",
            "message": "Native OMML inserted successfully.",
        }

    except Exception as exc:
        if fallback == "image":
            return _fallback_image(latex, slide, bbox, output_dir, slide_no, slot_name, str(exc))
        return {
            "status": "failed",
            "message": f"OMML insert failed: {exc}",
        }


def _fallback_image(
    latex: str, slide, bbox: list[float],
    output_dir: Path, slide_no: int, slot_name: str,
    error_msg: str = "",
) -> dict:
    """回退到 PNG 图片渲染（调用 render_formula.py）。"""
    from pptx.util import Inches

    formula_dir = output_dir / "assets" / "formulas"
    formula_dir.mkdir(parents=True, exist_ok=True)
    safe_slot = re.sub(r"[^a-zA-Z0-9_\-]+", "_", slot_name)
    out_path = formula_dir / f"slide_{slide_no:02d}_{safe_slot}.png"

    try:
        sys.path.insert(0, str(SCRIPT_DIR))
        from render_formula import render_formula
        render_formula(latex, out_path)
    except Exception:
        out_path = None

    if out_path and out_path.exists():
        x, y, w, h = [Inches(float(v)) for v in bbox]
        try:
            from PIL import Image
            with Image.open(out_path) as im:
                iw, ih = im.size
            box_ratio = float(w) / float(h)
            img_ratio = iw / ih
            if img_ratio > box_ratio:
                new_w = w
                new_h = int(float(w) / img_ratio)
            else:
                new_h = h
                new_w = int(float(h) * img_ratio)
            new_x = x + int((float(w) - float(new_w)) / 2)
            new_y = y + int((float(h) - float(new_h)) / 2)
        except Exception:
            new_x, new_y, new_w, new_h = x, y, w, h
        slide.shapes.add_picture(str(out_path), new_x, new_y, width=new_w, height=new_h)
        return {
            "status": "fallback_image",
            "image_path": str(out_path),
            "message": f"Native OMML failed ({error_msg}), used PNG fallback.",
        }

    return {
        "status": "failed",
        "message": f"Native OMML and PNG fallback both failed. OMML error: {error_msg}",
    }


# ── CLI (测试用) ──────────────────────────────────────────────────────


def main() -> int:
    """CLI 用于测试 LaTeX → OMML 转换流水线。"""
    if len(sys.argv) < 3:
        print("Usage: latex_to_omml.py <latex_string> <output_xml>")
        print("  Converts LaTeX math to OMML XML and saves to file.")
        print("  Use '-' for output_xml to print to stdout.")
        return 2

    latex = sys.argv[1]
    output = sys.argv[2]

    import lxml.etree as LX

    omml = latex_to_omml(latex)
    xml_str = LX.tostring(omml, pretty_print=True, encoding="unicode")

    if output == "-":
        print(xml_str)
    else:
        Path(output).parent.mkdir(parents=True, exist_ok=True)
        Path(output).write_text(xml_str, encoding="utf-8")
        print(f"Wrote OMML to {output}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
